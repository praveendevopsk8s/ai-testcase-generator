import streamlit as st
import subprocess
import pandas as pd
import PyPDF2
from docx import Document
from pptx import Presentation
import pyshark
import re

# --- Page Config ---
st.set_page_config(page_title="AI Test Case Generator", page_icon="🧪")

# --- Constants ---
MAX_FILE_SIZE = 200 * 1024 * 1024
OLLAMA_PATH = "ollama"  # Replace with full path if needed

# --- Load File Content ---
def load_file(file):
    try:
        if file.size > MAX_FILE_SIZE:
            st.error("File exceeds 200MB limit.")
            return None

        file_type = file.type
        if file_type in ["text/plain", "application/octet-stream"]:
            return file.read().decode("utf-8")
        elif file_type == "application/vnd.tcpdump.pcap":
            capture = pyshark.FileCapture(file)
            return "\n".join([str(packet) for packet in capture])
        elif file_type == "application/pdf":
            reader = PyPDF2.PdfReader(file)
            return ''.join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file)
            return '\n'.join([p.text for p in doc.paragraphs if p.text])
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            ppt = Presentation(file)
            return '\n'.join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text])
        else:
            st.error("Unsupported file type.")
            return None
    except Exception as e:
        st.error(f"Error reading file: {str(e)}")
        return None

# --- Ask LLaMA ---
def ask_llama(prompt, context=""):
    try:
        input_text = f"""
You are a QA engineer. Based on the following requirement document, generate test cases.

Context:
{context}

Instruction:
ONLY return the test cases in the following format:
<Actual test description> || <Steps to be followed> || <Expected Results>

Do NOT use labels like 'Test Case 1'. Just plain rows with || separator.
No headers, no markdown, no numbering, no explanations — only valid data rows.
"""
        result = subprocess.run(
            [OLLAMA_PATH, "run", "llama3.2"],
            input=input_text,
            text=True,
            capture_output=True,
            encoding="utf-8"
        )
        return result.stdout.strip()
    except Exception as e:
        return f"⚠️ Error: {str(e)}"

# --- Parse and Clean Output ---
def parse_test_cases(text):
    lines = [line.strip() for line in text.split("\n") if "||" in line]
    rows = []
    for line in lines:
        if re.match(r"(?i)test case\s*\d+\s*\|\|", line):
            line = re.sub(r"(?i)test case\s*\d+\s*\|\|", "", line)

        parts = [p.strip() for p in line.split("||")]
        if len(parts) == 3:
            rows.append(parts)

    df = pd.DataFrame(rows, columns=["Test Description", "Steps to be followed", "Expected Results"])
    return df

# --- Streamlit UI ---
st.markdown("## ✅ AI-Powered Test Case Generator")
st.markdown("Upload a requirements document and generate structured test cases.")

# Upload and save file content to session state
file = st.file_uploader("📄 Upload your document", type=["txt", "log", "csv", "pdf", "pptx", "docx", "pcap"])
if file:
    file_content = load_file(file)
    if file_content:
        st.session_state["file_content"] = file_content
        st.success("✅ Document uploaded successfully!")

# Prompt and generate button
user_prompt = st.text_input("✍️ Enter your prompt (e.g., 'Generate test cases for Wi-Fi 7')")

if st.button("🚀 Generate Test Cases"):
    if user_prompt and "file_content" in st.session_state:
        with st.spinner("Generating..."):
            response = ask_llama(user_prompt, context=st.session_state["file_content"])

            if not response or "⚠️" in response:
                st.error("❌ Failed to generate test cases. Try adjusting your prompt.")
            else:
                df = parse_test_cases(response)
                if df.empty:
                    st.warning("⚠️ No valid test cases parsed. Try a clearer prompt.")
                else:
                    st.markdown("### 📋 Test Case Table Preview")
                    st.dataframe(df, use_container_width=True)

                    csv = df.to_csv(index=False).encode("utf-8")
                    st.download_button(
                        label="📥 Download as CSV",
                        data=csv,
                        file_name="test_cases.csv",
                        mime="text/csv"
                    )
    else:
        st.warning("Please upload a file and enter a prompt before generating.")
